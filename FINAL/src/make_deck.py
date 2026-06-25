"""
Themed PowerPoint OUTLINE template: 'The Interlace x Graph Machine Learning'.
Visual language matched to the reference deck BioSpatial_Intelligence.pptx
(dark navy bg, purple #6C63FF accent, Montserrat + Consolas, corner brackets, dot-matrix).
Output: FINAL/out/The_Interlace_GraphML_Outline.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------------- theme: Concrete & Foliage (The Interlace) ----------------
# (constant names kept; PURPLE/PUR_* now hold the GREEN accent)
BG     = RGBColor(0x1E, 0x29, 0x22)   # deep forest charcoal
PANEL  = RGBColor(0x28, 0x35, 0x2D)   # card panel
PANEL2 = RGBColor(0x33, 0x45, 0x3A)   # lighter panel (big section number)
SLATE  = RGBColor(0x33, 0x45, 0x3A)   # corner brackets / dividers
LINE   = RGBColor(0x44, 0x57, 0x4B)   # thin lines
PURPLE = RGBColor(0x6F, 0xA3, 0x6B)   # primary accent (foliage green)
PUR_L  = RGBColor(0xA7, 0xC4, 0xA0)   # light green
LILAC  = RGBColor(0xC4, 0xD6, 0xBE)   # pale green
GREY   = RGBColor(0x8A, 0x98, 0x8A)   # muted text
GREYL  = RGBColor(0xB6, 0xC4, 0xB2)   # light muted
TLIGHT = RGBColor(0xE8, 0xEE, 0xE8)   # body text on dark
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DTEXT  = RGBColor(0x23, 0x2A, 0x25)   # text on light
# --- light theme (content slides) ---
PAPER  = RGBColor(0xF2, 0xF4, 0xEE)   # warm off-white
CARD_L = RGBColor(0xE6, 0xED, 0xE2)   # light green card
BRK_L  = RGBColor(0xCB, 0xD8, 0xC4)   # light corner brackets
PUR_D  = RGBColor(0x4C, 0x7A, 0x48)   # deeper green for text on light
SUB_D  = RGBColor(0x5E, 0x6B, 0x5C)   # muted text on light

HEAD = "Montserrat"
MONO = "Consolas"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

GROUP = "G05"   # <- change if needed


def _no_line(sp): sp.line.fill.background()
def _no_shadow(sp): sp.shadow.inherit = False

def rect(slide, x, y, w, h, color, line=None, line_w=None, dash=False, rounded=False):
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sp = slide.shapes.add_shape(shp, x, y, w, h)
    if color is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        _no_line(sp)
    else:
        sp.line.color.rgb = line; sp.line.width = line_w or Pt(1)
        if dash:
            ln = sp.line._get_or_add_ln()
            ln.append(ln.makeelement(qn('a:prstDash'), {'val': 'dash'}))
    _no_shadow(sp)
    return sp

def txt(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=1.0):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = space
        for (t, sz, col, fn, bold) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col
            r.font.name = fn; r.font.bold = bold
    return tb

def bullets(slide, x, y, w, h, items, size=16, color=TLIGHT, gap=8, mark_color=PURPLE, sub_color=GREYL):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, it in enumerate(items):
        lvl, text = (it if isinstance(it, tuple) else (0, it))
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.line_spacing = 1.08; p.level = lvl
        mk = p.add_run(); mk.text = ("▪  " if lvl == 0 else "·  ")
        mk.font.size = Pt(size if lvl == 0 else size - 1)
        mk.font.color.rgb = mark_color if lvl == 0 else sub_color
        mk.font.name = MONO; mk.font.bold = False
        r = p.add_run(); r.text = text
        r.font.size = Pt(size if lvl == 0 else size - 1)
        r.font.color.rgb = color if lvl == 0 else sub_color
        r.font.name = HEAD; r.font.bold = False
    return tb

def corner_brackets(slide, color=SLATE):
    a = Inches(0.32); th = Inches(0.022)
    pts = [(Inches(0.32), Inches(0.32)),
           (SW-Inches(0.32)-a, Inches(0.32)),
           (Inches(0.32), SH-Inches(0.32)-a),
           (SW-Inches(0.32)-a, SH-Inches(0.32)-a)]
    for (x, y) in pts:
        rect(slide, x, y, a, th, color)
        rect(slide, x, y, th, a, color)
        rect(slide, x+a-th, y, th, a, color)
        rect(slide, x, y+a-th, a, th, color)

def dot_matrix(slide, x, y, cols=6, rows=2, color=PURPLE):
    d = Inches(0.05); gap = Inches(0.22)
    for r in range(rows):
        for c in range(cols):
            sq = rect(slide, Emu(int(x+c*gap)), Emu(int(y+r*gap)), d, d, color)

def graph_motif(slide, cx, cy, scale=1.0, node=PURPLE, edge=LINE):
    pts = [(0.0,0.0),(1.2,-0.5),(2.1,0.6),(1.0,1.3),(-0.2,1.0),(2.6,1.6),(0.6,-1.1)]
    edges = [(0,1),(1,2),(2,3),(3,4),(4,0),(2,5),(3,5),(0,6)]
    def P(p): return (Emu(int(cx+p[0]*Inches(0.85)*scale)), Emu(int(cy+p[1]*Inches(0.85)*scale)))
    for a,b in edges:
        x1,y1=P(pts[a]); x2,y2=P(pts[b])
        ln=slide.shapes.add_connector(2,x1,y1,x2,y2)
        ln.line.color.rgb=edge; ln.line.width=Pt(1.25)
    rr=Inches(0.13)*scale
    for p in pts:
        x,y=P(p)
        d=slide.shapes.add_shape(MSO_SHAPE.OVAL,Emu(int(x-rr)),Emu(int(y-rr)),Emu(int(2*rr)),Emu(int(2*rr)))
        d.fill.solid(); d.fill.fore_color.rgb=node; _no_line(d); _no_shadow(d)

def base(slide, bg=BG, brk=SLATE):
    rect(slide, 0, 0, SW, SH, bg)
    corner_brackets(slide, brk)

def footer(slide, n):
    txt(slide, Inches(0.55), Inches(7.02), Inches(7), Inches(0.3),
        [[("THE INTERLACE  ×  GRAPH ML", 9, GREY, MONO, False)]])
    txt(slide, Inches(11.5), Inches(7.02), Inches(1.3), Inches(0.3),
        [[(f"{n:02d} / {GROUP}", 9, GREY, MONO, False)]], align=PP_ALIGN.RIGHT)

def placeholder(slide, x, y, w, h, label, light=True):
    fill = CARD_L if light else PANEL
    txtc = PUR_D if light else PUR_L
    card = rect(slide, x, y, w, h, fill, line=PURPLE, line_w=Pt(1.25), dash=True, rounded=True)
    tf = card.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "▤  " + label
    r.font.size = Pt(13); r.font.color.rgb = txtc; r.font.name = MONO; r.font.bold = False
    return card

_pageno = [0]

def content(title, kicker, items=None, ph=None, ph2=None, two_col=None):
    _pageno[0] += 1
    s = prs.slides.add_slide(BLANK); base(s, PAPER, BRK_L)
    txt(s, Inches(0.78), Inches(0.62), Inches(11), Inches(0.3),
        [[(kicker.upper(), 11, PUR_D, MONO, False)]])
    txt(s, Inches(0.75), Inches(0.95), Inches(12), Inches(0.95),
        [[(title, 30, DTEXT, HEAD, True)]])
    rect(s, Inches(0.78), Inches(1.72), Inches(0.7), Pt(4), PURPLE)
    by = Inches(2.15)
    if two_col:
        left, right = two_col
        bullets(s, Inches(0.78), by, Inches(5.7), Inches(4.5), left, color=DTEXT, sub_color=SUB_D)
        bullets(s, Inches(7.0), by, Inches(5.5), Inches(4.5), right, color=DTEXT, sub_color=SUB_D)
    elif items and (ph or ph2):
        bullets(s, Inches(0.78), by, Inches(5.7), Inches(4.6), items, color=DTEXT, sub_color=SUB_D)
        if ph:  placeholder(s, Inches(6.95), by, Inches(5.6), Inches(2.1), ph)
        if ph2: placeholder(s, Inches(6.95), Inches(4.5), Inches(5.6), Inches(2.1), ph2)
    elif items:
        bullets(s, Inches(0.78), by, Inches(11.8), Inches(4.6), items, size=17, color=DTEXT, sub_color=SUB_D)
    elif ph:
        placeholder(s, Inches(0.78), by, Inches(11.8), Inches(4.55), ph)
    footer(s, _pageno[0])
    return s

def section(num, title, subtitle=""):
    _pageno[0] += 1
    s = prs.slides.add_slide(BLANK); base(s)
    txt(s, Inches(0.7), Inches(1.7), Inches(6), Inches(2.4),
        [[(f"{num:02d}", 150, PANEL2, HEAD, True)]])
    rect(s, Inches(0.82), Inches(4.05), Inches(0.9), Pt(5), PURPLE)
    txt(s, Inches(0.8), Inches(4.2), Inches(11), Inches(1.0),
        [[(title, 34, WHITE, HEAD, True)]])
    if subtitle:
        txt(s, Inches(0.82), Inches(5.1), Inches(11), Inches(0.6),
            [[(subtitle, 15, GREYL, HEAD, False)]])
    dot_matrix(s, SW-Inches(1.92), SH-Inches(1.2), cols=6, rows=2, color=PANEL2)
    footer(s, _pageno[0])
    return s

# ==================== TITLE ====================
_pageno[0] += 1
s = prs.slides.add_slide(BLANK); base(s)
dot_matrix(s, Inches(0.78), Inches(0.78))
txt(s, Inches(0.78), Inches(1.32), Inches(11.5), Inches(0.3),
    [[("IAAC · MaCAD  |  Digital tools for Graph Machine Learning", 12, GREYL, MONO, False)]])
txt(s, Inches(11.2), Inches(1.32), Inches(1.6), Inches(0.3),
    [[("BUILDINGS AS GRAPHS", 9, GREY, MONO, False)]], align=PP_ALIGN.RIGHT)
rect(s, Inches(0.8), Inches(2.6), Inches(0.9), Pt(5), PURPLE)
txt(s, Inches(0.78), Inches(2.15), Inches(6), Inches(0.3),
    [[("FINAL ASSIGNMENT · " + GROUP, 12, PUR_L, MONO, False)]])
txt(s, Inches(0.75), Inches(2.78), Inches(11.6), Inches(1.7),
    [[("The Interlace", 56, WHITE, HEAD, True)],
     [("Graph Machine Learning", 30, PUR_L, HEAD, False)]], space=1.0)
txt(s, Inches(0.78), Inches(4.95), Inches(11.0), Inches(0.5),
    [[("A graph-based analysis & node-classification pipeline for architectural floor plans", 15, GREYL, HEAD, False)]])
txt(s, Inches(0.78), Inches(5.5), Inches(11.8), Inches(0.3),
    [[("Nihan Malkoç  ·  Sushmitha Ravi  ·  Tue Minh (Chloe)  ·  Zeynep Sezen Dursun", 13, TLIGHT, HEAD, False)]])
rect(s, Inches(0.8), Inches(6.0), Inches(0.55), Pt(3), LINE)
txt(s, Inches(0.78), Inches(6.15), Inches(11.5), Inches(0.6),
    [[("Faculties:  ", 12, GREY, MONO, False),
      ("Wassim Jabi  |  Olga Poletkina", 12, GREYL, HEAD, False)]])
txt(s, Inches(0.78), Inches(6.5), Inches(11.5), Inches(0.4),
    [[("Dataset: Modified Swiss Dwellings (MSD)   ·   Tools: TopologicPy + PyTorch Geometric", 10, GREY, MONO, False)]])
graph_motif(s, Inches(9.7), Inches(2.7), scale=1.05, node=PURPLE, edge=LINE)
footer(s, _pageno[0])

# ==================== AGENDA ====================
content("Agenda", "Overview", items=[
    "01  Introduction & Objective",
    "02  Dataset & Reference Research",
    "03  Methodology — Floor Plan to Graph",
    "04  Spatial / Graph Analysis",
    "05  Node Classification",
    "06  Discussion & Conclusion",
])

# ===== SECTION 1 =====
section(1, "Introduction & Objective", "What we build and why The Interlace")
content("Project Objective", "Introduction", items=[
    "Build a graph-based analysis & learning pipeline for architectural floor plans.",
    "Recreate a multi-apartment floor plan, represent it as a graph, analyze spatial organization.",
    "Run node classification (predict room type) with a GraphSAGE model.",
    "Goal: understanding & interpretation — not perfect model accuracy.",
    "Case study: The Interlace (Singapore) — floors 3 & 4.",
])
content("Why The Interlace?", "Introduction",
    items=[
        "Multi-apartment residential complex — many units in clusters (mirrors MSD).",
        "Interconnected blocks + courtyards → rich spatial relationships.",
        "Strong circulation: shared corridors, vertical cores, communal spaces.",
        "Diverse space types → ideal for node classification & centrality analysis.",
    ],
    ph="IMAGE: The Interlace photo / massing")

# ===== SECTION 2 =====
section(2, "Dataset & Reference Research", "Modified Swiss Dwellings + papers")
content("Modified Swiss Dwellings (MSD)", "Dataset",
    items=[
        "~5,372 floor plans of medium–large building complexes; 18.9k apartments.",
        "Each plan in 3 modalities: raster image, vector geometry, graph.",
        "Areas (rooms), separators (walls), openings (doors/windows).",
        "Node attributes: centroid, geometry (polygon), roomtype, zonetype.",
        "Connectivity types: passage, door, front door.",
    ],
    ph="IMAGE: MSD sample plan + access graph")
content("Reference Papers", "Dataset", two_col=(
    [ "Paper 1 — MSD: A Benchmark Dataset for Floor Plan Generation of Building Complexes",
      (1, "Multi-apartment focus; access-graph representation"),
      (1, "Room-type vs zone-type categories"),
      (1, "Standard splits for benchmarking"),
    ],
    [ "Paper 2 — GNNs for Node Classification & Attribute Allocation in Architectural BIM",
      (1, "Rooms as nodes, adjacency as edges"),
      (1, "GNN predicts room / element attributes"),
      (1, "Motivates our GraphSAGE node-classifier"),
    ]))
content("How Floor Plans Are Represented", "Dataset",
    items=[
        "Node = room (area). Position = polygon centroid.",
        "9 room-type classes: Bedroom, Bathroom, Corridor, Balcony, Kitchen, Livingroom, Stairs, Storeroom, Dining.",
        "4 zone types: private (Z1), public (Z2), service (Z3), outside (Z4).",
        "Edge = adjacency / access between rooms.",
        "Graph attributes feed graph-based learning.",
    ],
    ph="DIAGRAM: plan → nodes(rooms) + edges(adjacency)")

# ===== SECTION 3 =====
section(3, "Methodology", "From floor plan to graph with TopologicPy")
content("Pipeline Overview", "Methodology",
    ph="DIAGRAM: Rhino (OBJ per layer) → TopologicPy graph → nodes/edges/graphs.csv → PyG GraphSAGE → prediction")
content("Recreated Floor Plan — The Interlace", "Methodology",
    items=[
        "Floors 3 & 4 recreated in Rhino, following MSD logic (not an exact copy).",
        "Each room type on its own layer → exported separately as OBJ.",
        "Rooms, apartments, corridors, cores, balconies, courtyards.",
    ],
    ph="OUTPUT: recreated floor plan (Rhino)")
content("Graph Construction (TopologicPy)", "Methodology",
    items=[
        "Import OBJ layers → tag each room with its type.",
        "Rooms → nodes (centroid); shared walls → edges (adjacency).",
        "Same node/edge schema & attributes as the MSD reference.",
        "Export to graphs.csv / nodes.csv / edges.csv.",
    ],
    ph="OUTPUT: TopologicPy graph over the plan")
content("Node Features & Class Schema", "Methodology", two_col=(
    [ "Node features (geometry, no label leakage):",
      (1, "feat_0  area"),
      (1, "feat_1  perimeter"),
      (1, "feat_2  compactness"),
      (1, "feat_3  number of corners"),
      (1, "feat_4  bounding-box aspect"),
      (1, "feat_5  rectangularity"),
    ],
    [ "Target label — roomtype (0–8):",
      (1, "0 Bedroom    1 Bathroom"),
      (1, "2 Corridor   3 Balcony"),
      (1, "4 Kitchen    5 Livingroom"),
      (1, "6 Stairs     7 Storeroom"),
      (1, "8 Dining"),
      "+ Layer → class mapping table (Interlace).",
    ]))

# ===== SECTION 4 =====
section(4, "Spatial / Graph Analysis", "What the metrics reveal")
content("Analysis Overview", "Analysis",
    items=[
        "Compute classic graph metrics on the floor-plan graph (as in Assignment 2).",
        "Degree centrality · Closeness centrality · Clustering coefficient · Shortest paths.",
        "Interpret: accessibility, circulation, apartment organization, spatial hierarchy, connectivity.",
    ])
content("Degree Centrality", "Analysis",
    items=["Highly connected rooms = circulation hubs (corridors, lobbies).", "Interpretation: [fill in]"],
    ph="OUTPUT: degree-centrality graph map")
content("Closeness Centrality", "Analysis",
    items=["How easily a room reaches all others → accessibility.", "Interpretation: [fill in]"],
    ph="OUTPUT: closeness-centrality graph map")
content("Clustering Coefficient", "Analysis",
    items=["Local grouping → apartment clusters / units.", "Interpretation: [fill in]"],
    ph="OUTPUT: clustering visualization")
content("Shortest Paths", "Analysis",
    items=["Circulation routes & distances between units / cores.", "Interpretation: [fill in]"],
    ph="OUTPUT: shortest-path example")
content("Interpretation — Spatial Reading", "Analysis", two_col=(
    [ "Accessibility & circulation:",
      (1, "[which rooms are hubs]"),
      (1, "[core / corridor role]"),
      "Apartment organization:",
      (1, "[clusters per unit]"),
    ],
    [ "Spatial hierarchy:",
      (1, "[public vs private]"),
      "Connectivity between units:",
      (1, "[how units link via shared spaces]"),
    ]))

# ===== SECTION 5 =====
section(5, "Node Classification", "Predicting room types with GraphSAGE")
content("Node Classification Pipeline", "Classification",
    items=[
        "Model: GraphSAGE node classifier (topologicpy.PyG), trained on MSD.",
        "Input: node features + graph adjacency; output: room-type class (0–8).",
        "Apply trained model to The Interlace graph → predict each room.",
        "Evaluate predictions, not chase perfect accuracy.",
    ],
    ph="DIAGRAM: GraphSAGE node-classification")
content("Training Results (MSD)", "Classification",
    items=["Train / val / test on MSD floor-plan graphs.",
           "Report accuracy / per-class metrics.",
           "Confusion matrix → which room types get confused."],
    ph="OUTPUT: training curves", ph2="OUTPUT: confusion matrix")
content("Predictions on The Interlace", "Classification",
    items=["Predicted vs actual room type per node.",
           "Correct vs incorrect predictions highlighted."],
    ph="OUTPUT: predicted graph", ph2="OUTPUT: actual (ground-truth) graph")
content("Do Predictions Make Spatial Sense?", "Classification",
    items=[
        "Where does the model succeed / fail, and why?",
        "Are errors spatially reasonable (e.g., bedroom ↔ studio)?",
        "Role of geometry vs neighborhood in the prediction.",
    ],
    ph="OUTPUT: error-analysis highlight")

# ===== SECTION 6 =====
section(6, "Discussion & Conclusion", "Limitations, improvements, takeaways")
content("Limitations & Improvements", "Discussion", two_col=(
    [ "Limitations:",
      (1, "Domain gap: MSD (CH dwellings) vs The Interlace"),
      (1, "Recreated plan ≠ exact building"),
      (1, "Class imbalance (e.g., Dining rare)"),
      (1, "Adjacency definition simplifications"),
    ],
    [ "Improvements:",
      (1, "More / richer node features"),
      (1, "Fine-tune on similar typologies"),
      (1, "Add door / access edges, not just adjacency"),
      (1, "Multi-floor (3D) graph"),
    ]))
content("Conclusion", "Discussion", items=[
    "Graph representation makes spatial structure explicit and analyzable.",
    "Centrality & clustering reveal circulation and apartment organization.",
    "GraphSAGE can infer room types from geometry + topology.",
    "The Interlace is a strong, legible case study for graph learning.",
])

# ==================== THANK YOU ====================
_pageno[0] += 1
s = prs.slides.add_slide(BLANK); base(s)
dot_matrix(s, Inches(0.78), Inches(0.78))
rect(s, Inches(0.8), Inches(3.0), Inches(0.9), Pt(5), PURPLE)
txt(s, Inches(0.78), Inches(3.15), Inches(11), Inches(1.1),
    [[("Thank you", 46, WHITE, HEAD, True)]])
txt(s, Inches(0.8), Inches(4.3), Inches(11), Inches(0.5),
    [[("Questions & discussion", 16, PUR_L, HEAD, False)]])
txt(s, Inches(0.8), Inches(4.9), Inches(11.5), Inches(0.5),
    [[("References: MSD (van Engelenburg et al., ECCV 2024) · GNNs for Node Classification in BIM · TopologicPy · PyTorch Geometric", 10, GREY, MONO, False)]])
graph_motif(s, Inches(9.7), Inches(2.4), scale=1.05, node=PURPLE, edge=LINE)
footer(s, _pageno[0])

out = r"C:/PROJECTS/GRAPH-ML-DOCUMENTS/FINAL/out/The_Interlace_GraphML_Outline.pptx"
prs.save(out)
print("Saved:", out)
print("Slides:", len(prs.slides._sldIdLst))
